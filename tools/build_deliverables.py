"""Generate docs/ECE_6255_Report.docx and docs/ECE_6255_Term_Project_Video_Slides.pptx."""
from pathlib import Path

from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.shared import Inches, Pt
from pptx import Presentation
from pptx.util import Inches as PptInches, Pt as PptPt

ROOT = Path(__file__).resolve().parents[1]
DOCS = ROOT / "docs"
REPO_URL = "https://github.com/TopelJL/ECE-6255-Digital-Processing-Speech-Signals-Term-Project"


def build_docx():
    DOCS.mkdir(parents=True, exist_ok=True)
    d = Document()
    t = d.add_heading("ECE 6255 — Term Project Report", 0)
    t.alignment = WD_ALIGN_PARAGRAPH.CENTER

    d.add_paragraph("Digital Processing of Speech Signals — Speech segmentation (voiced / unvoiced / silence)")
    d.add_paragraph("Georgia Institute of Technology, Spring 2026")
    d.add_paragraph("Instructor: B.H. Juang")
    d.add_paragraph("Team: Jaxon Topel, Adrian Cruz, Michael Ritz")

    d.add_heading("1.0 Project summary", level=1)
    d.add_paragraph(
        "We segment speech into three classes per short-time frame: silence, voiced, and unvoiced, "
        "using short-time energy (STE), zero-crossing rate (ZCR), and autocorrelation-based pitch strength, "
        "followed by median filtering and optional higher-level boundary cues."
    )

    d.add_heading("1.1 Notes for consistency", level=1)
    d.add_paragraph(
        "Proofreading: use the heading spelling silence (not “slicence”) in §1.1.1. "
        "Use “Michael Ritz” consistently (not “Ruiz”) across report and slides."
    )

    d.add_heading("1.2 Theoretical foundation (equations)", level=1)
    d.add_paragraph(
        "Let x_w[n], n = 0,…,N−1, be a windowed frame (Hamming window). Short-time energy (STE):"
    )
    p = d.add_paragraph()
    p.add_run("E = (1/N) · Σₙ x_w[n]²").italic = True

    d.add_paragraph("Zero-crossing rate (ZCR) for the windowed frame:")
    p = d.add_paragraph()
    p.add_run("ZCR = (1/(2N)) · Σₙ |sgn(x_w[n]) − sgn(x_w[n−1])|").italic = True
    d.add_paragraph("(with sgn(0) treated as +1 to avoid undefined crossings).")

    d.add_paragraph("Short-time autocorrelation at lag k (unnormalized form):")
    p = d.add_paragraph()
    p.add_run("R(k) = Σₙ x_w[n] · x_w[n+k]").italic = True
    d.add_paragraph(
        "Voicing uses the normalized autocorrelation (coefficient scaling). A prominent peak in the "
        "physiological pitch lag range implies periodicity → voiced speech."
    )

    d.add_heading("1.3 / 2.1 Adaptive thresholds", level=1)
    d.add_paragraph(
        "Silence is declared when normalized STE falls below a low percentile of the utterance’s STE distribution. "
        "For non-silent frames, a ZCR cutoff is derived from the distribution of ZCR on high-energy (speech-like) frames. "
        "Voiced frames require sufficiently high autocorrelation peak height in the F0 lag band and ZCR below that adaptive cutoff; "
        "otherwise the frame is unvoiced. Parameters live in params.m (quantiles and voicing_peak_min)."
    )

    d.add_heading("3.0 Source code", level=1)
    d.add_paragraph(f"Public repository: {REPO_URL}")

    d.add_heading("3.1 Instructions", level=1)
    d.add_paragraph("1. Install MATLAB (R2019b+ recommended).")
    d.add_paragraph("2. git clone the repository URL above.")
    d.add_paragraph("3. Open MATLAB and cd to the project root (folder containing main.m).")
    d.add_paragraph("4. Run: main")
    d.add_paragraph("5. Select a .wav file when prompted (example_wavs/example1.wav or example2.wav).")
    d.add_paragraph("6. Inspect the figure and files under results/ (*.png, *.fig).")
    d.add_paragraph("7. For metrics on the synthetic clip: run evaluate_example2 in MATLAB.")

    d.add_heading("3.2 Code running example (video)", level=1)
    d.add_paragraph(
        "[INSERT LINK after upload — e.g. unlisted YouTube URL for the 12-minute team video.]"
    )

    d.add_heading("4.0–4.1 Results (fill after running pipeline)", level=1)
    d.add_paragraph(
        "Run evaluate_example2.m and paste accuracy and the 3×3 confusion matrix from results/metrics_example2.txt. "
        "Add waveform + segmentation figures from results/ for example1 and example2 (and any team recordings)."
    )
    d.add_paragraph("Example table placeholder:")
    table = d.add_table(rows=4, cols=4)
    hdr = table.rows[0].cells
    hdr[0].text = "Metric"
    hdr[1].text = "Value"
    hdr[2].text = "Notes"
    hdr[3].text = ""
    table.rows[1].cells[0].text = "Frame accuracy (example2)"
    table.rows[1].cells[1].text = "Run MATLAB"
    table.rows[2].cells[0].text = "Confusion matrix"
    table.rows[2].cells[1].text = "See results/"
    table.rows[3].cells[0].text = "Qualitative (example1)"
    table.rows[3].cells[1].text = "Inspect plots"

    d.add_heading("5.0 Conclusion", level=1)
    d.add_paragraph(
        "Classical features separate silence, voiced, and unvoiced segments without ML; adaptive thresholds "
        "improve robustness across level variations. Future work: spectral flux / centroid, TIMIT alignment, phoneme-level boundaries."
    )

    d.add_heading("6.0 References", level=1)
    refs = [
        "L. R. Rabiner and R. W. Schafer, Theory and Applications of Digital Speech Processing, Pearson, 2011.",
        "MathWorks documentation: audioread, hamming, fft, xcorr, medfilt1 — as used in the codebase.",
        "Course materials: ECE 6255, Georgia Tech (Spring 2026), B.H. Juang.",
    ]
    for r in refs:
        d.add_paragraph(r, style="List Bullet")

    out = DOCS / "ECE_6255_-_Term_Project.docx"
    d.save(out)
    print("Wrote", out)


def build_pptx():
    DOCS.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = PptInches(13.333)
    prs.slide_height = PptInches(7.5)

    def add_title_slide(title, subtitle):
        layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = title
        slide.placeholders[1].text = subtitle

    def add_bullet_slide(title, bullets):
        layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = title
        body = slide.placeholders[1].text_frame
        body.clear()
        for i, b in enumerate(bullets):
            p = body.paragraphs[0] if i == 0 else body.add_paragraph()
            p.text = b
            p.level = 0
            p.font.size = PptPt(20)

    add_title_slide(
        "Speech segmentation: Voiced / Unvoiced / Silence",
        "ECE 6255 — B.H. Juang — Team: Topel, Cruz, Ritz",
    )

    add_bullet_slide(
        "Project overview",
        [
            "Goal: label each short-time frame as silence, voiced, or unvoiced.",
            "Why: front-end for recognition, coding, and analysis.",
            "Approach: STE + ZCR + autocorrelation pitch strength; classical DSP only.",
        ],
    )

    add_bullet_slide(
        "Theoretical foundation",
        [
            "STE E = (1/N) Σ x_w² — high for voiced, very low for silence.",
            "ZCR — high for noise-like unvoiced fricatives; lower for periodic voiced.",
            "R(k) peak in pitch lag range → periodicity → voiced.",
        ],
    )

    add_bullet_slide(
        "Algorithm pipeline",
        [
            "Frame: ~25 ms, Hamming, ~10 ms hop.",
            "Features per frame → adaptive thresholds → 3-class labels.",
            "Median filter + min run merge → optional boundary heuristics.",
            "Plots + export to results/.",
        ],
    )

    add_bullet_slide(
        "Implementation (MATLAB)",
        [
            "main.m — entry, file dialog.",
            "compute_features.m — STE, ZCR, pitch strength.",
            "classify_frames.m + smooth_labels.m — decisions + medfilt1.",
            "params.m — all tunables in one place.",
        ],
    )

    add_bullet_slide(
        "Adaptive thresholds",
        [
            "Silence: STE below a low percentile of the utterance.",
            "Speech mask from STE quantile drives ZCR statistics.",
            "Voiced: high autocorr peak + ZCR below adaptive cutoff.",
        ],
    )

    add_bullet_slide(
        "Results — plots",
        [
            "Insert figures from results/ after running main.m.",
            "Show example1 (speech) and example2 (synthetic regions).",
        ],
    )

    add_bullet_slide(
        "Results — metrics",
        [
            "Run evaluate_example2.m.",
            "Report frame accuracy and 3×3 confusion matrix.",
            "Optional: hand labels for more utterances (see evaluation/README.md).",
        ],
    )

    add_bullet_slide(
        "Code demo",
        [
            "Screen recording: MATLAB → main → pick example_wavs.",
            "Full walkthrough in the 12-minute team video.",
        ],
    )

    add_bullet_slide(
        "Conclusion & future work",
        [
            "Median smoothing reduces single-frame errors.",
            "Future: spectral flux/centroid, TIMIT comparison, phoneme boundaries.",
        ],
    )

    add_bullet_slide(
        "References",
        [
            "Rabiner & Schafer, Digital Speech Processing, 2011.",
            "ECE 6255 course materials (Juang).",
            f"Code: {REPO_URL}",
        ],
    )

    out = DOCS / "ECE_6255_-_Term_Project_Video_Slides.pptx"
    prs.save(out)
    print("Wrote", out)


if __name__ == "__main__":
    build_docx()
    build_pptx()
